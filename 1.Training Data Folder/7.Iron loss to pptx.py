import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import io

# python-pptxライブラリをインポートします (pip install python-pptx が必要)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("【エラー】python-pptxライブラリが見つかりません。")
    print("コマンドプロンプトで `pip install python-pptx` を実行してインストールしてください。")
    exit()

# ==============================================================================
# ユーザー設定箇所
# ==============================================================================
# --- 材料名 ---
MAT_NAME = "50A470"  # 材料名を指定

# --- 周波数リスト ---
FREQ_LIST = [20, 50, 100, 200, 400, 600, 800, 1000]

# --- 処理する振幅の範囲 ---
AMP_START = 0.05
AMP_END = 1.8
AMP_STEP = 0.05

# --- グラフ表示設定 ---
PLOT_LINEAR_FIT = False

# ★★★ 修正箇所：相対パス設定 ★★★
# ----------------------------------------------------------------------
# 1. このスクリプト(.py)がある GPR_perf フォルダの場所を取得
base_dir = os.path.dirname(__file__)

# 2. base_dir からの相対的な位置で各パスを定義
#    「1.Training Data Folder」の重複がないように修正しました
BASE_PATH = os.path.join(base_dir)
INPUT_BASE_PATH = os.path.join(BASE_PATH, "3.Fourier Transform Correction")
OUTPUT_BASE_PATH = os.path.join(BASE_PATH, "8.Hysteresis Loss")
# ----------------------------------------------------------------------


# ==============================================================================
# プログラム本体
# ==============================================================================

def find_input_file(directory, amp, freq):
    """
    指定された振幅の入力ファイルを探す。
    """
    if round(amp, 2) * 10 == int(round(amp, 2) * 10):
        amp_str = f"{amp:.1f}"
    else:
        amp_str = f"{amp:.2f}"
    
    filename = f"Bm{amp_str}hys_{freq}hz.xlsx"
    filepath = os.path.join(directory, filename)
    
    return filepath if os.path.exists(filepath) else None


def calculate_hysteresis_area(h_data, b_data):
    """シューレースの公式を用いてヒステリシスループの面積を計算する。"""
    return 0.5 * np.abs(np.dot(h_data, np.roll(b_data, -1)) - np.dot(b_data, np.roll(h_data, -1)))

def main():
    """メイン処理を実行する関数"""
    print(f"■ Material: {MAT_NAME}, Hysteresis Loss Calculation Start.")
    
    os.makedirs(OUTPUT_BASE_PATH, exist_ok=True)

    amp_list = np.round(np.arange(AMP_START, AMP_END + 1e-8, AMP_STEP), 2)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Excelのファイル名に周波数リストを含めると長くなりすぎるため、材料名のみにする
    excel_output_path = os.path.join(OUTPUT_BASE_PATH, f"Hysteresis_Analysis_Summary_{MAT_NAME}.xlsx")
    
    with pd.ExcelWriter(excel_output_path, engine='openpyxl') as writer:
        for amp in amp_list:
            print("\n" + "="*60)
            print(f"Processing for AMP = {amp:.2f} T")
            print("="*60)

            results_data = []
            for freq in FREQ_LIST:
                print(f"   - Processing Frequency: {freq} Hz...")
                
                input_dir = os.path.join(INPUT_BASE_PATH, MAT_NAME, f"{freq}Hz")
                file_path = find_input_file(input_dir, amp, freq)

                if file_path is None:
                    print(f"     - WARNING: File for AMP={amp:.2f}T, FREQ={freq}Hz not found. Skipping.")
                    results_data.append({'Frequency [Hz]': freq, 'Energy Loss [J/m^3]': np.nan})
                    continue

                try:
                    df = pd.read_excel(file_path, header=None, names=['H', 'B'])
                    h_values = df['H'].values
                    b_values = df['B'].values
                except Exception as e:
                    print(f"     - WARNING: Could not read file {os.path.basename(file_path)}. Skipping. Error: {e}")
                    results_data.append({'Frequency [Hz]': freq, 'Energy Loss [J/m^3]': np.nan})
                    continue
                
                energy_loss = calculate_hysteresis_area(h_values, b_values)
                results_data.append({'Frequency [Hz]': freq, 'Energy Loss [J/m^3]': energy_loss})
            
            loss_df = pd.DataFrame(results_data).dropna()
            if len(loss_df) < 3:
                print(f"   - Less than 3 valid data points for AMP={amp:.2f}T. Skipping fitting and plotting.")
                continue
                
            frequencies = loss_df['Frequency [Hz]'].values
            areas = loss_df['Energy Loss [J/m^3]'].values

            slope, intercept_linear = np.polyfit(frequencies, areas, 1)
            A = np.vstack([frequencies, np.sqrt(frequencies), np.ones(len(frequencies))]).T
            p, _, _, _ = np.linalg.lstsq(A, areas, rcond=None)
            a1, a2, intercept_3term = p

            fig_fit, ax = plt.subplots(figsize=(10, 7.5))
            
            ax.scatter(frequencies, areas, color='blue', zorder=5, label='Calculated Data')
            
            if PLOT_LINEAR_FIT:
                plot_x_linear = np.array([0, frequencies.max()])
                plot_y_linear = slope * plot_x_linear + intercept_linear
                label_linear = (f'Linear Fit\n'fr'$y = {slope:.4f}x + {intercept_linear:.4f}$')
                ax.plot(plot_x_linear, plot_y_linear, color='red', linestyle='--', label=label_linear)
                ax.scatter(0, intercept_linear, color='red', marker='x', s=120, zorder=10, label=fr'Intercept (c) = ${intercept_linear:.4f}$')
            
            plot_x_3term = np.linspace(0, frequencies.max(), 200)
            plot_y_3term = a1 * plot_x_3term + a2 * np.sqrt(plot_x_3term) + intercept_3term
            label_3term = (f'3-term Fit\n'fr'$y = {a1:.4f}x + {a2:.4f}\sqrt{{x}} + {intercept_3term:.4f}$')
            ax.plot(plot_x_3term, plot_y_3term, color='green', linestyle=':', label=label_3term)
            ax.scatter(0, intercept_3term, color='green', marker='o', s=100, zorder=10, facecolors='none', edgecolors='green', linewidth=2, label=fr'Intercept (b) = ${intercept_3term:.4f}$')
            
            ax.set_xlabel(fr'$\it f$ [Hz]', fontsize=18)
            ax.set_ylabel(fr'$\it P_i/f$ [J/m$^3$]', fontsize=18)
            ax.grid(True)
            ax.set_xlim(left=0)
            ax.set_ylim(bottom=0)
            ax.legend(fontsize=12)
            plt.tight_layout(rect=[0, 0.1, 1, 1])

            fig_num = int(round(amp / AMP_STEP))
            title_str = fr'Fig. {fig_num}. $B_{{\mathrm{{m}}}}$ = {amp:.2f} T'
            fig_fit.text(0.5, 0.02, title_str, ha='center', fontsize=20)
            
            image_stream = io.BytesIO()
            plt.savefig(image_stream, format='png', dpi=300)
            plt.close(fig_fit)
            
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.add_picture(image_stream, Inches(1.666), Inches(0), width=Inches(10.0))
            image_stream.close()
            print(f"   - Graph for AMP={amp:.2f}T has been added to PowerPoint.")

            fit_results_df = pd.DataFrame({
                'Parameter': ['Linear_Slope(m)', 'Linear_Intercept(c)', '3-term_Coeff(a1)', '3-term_Coeff(a2)', '3-term_Intercept(b)'],
                'Value': [slope, intercept_linear, a1, a2, intercept_3term]
            })
            sheet_name_loss = f'{amp:.2f}T Loss'; sheet_name_fit = f'{amp:.2f}T Fit'
            loss_df.to_excel(writer, sheet_name=sheet_name_loss, index=False)
            fit_results_df.to_excel(writer, sheet_name=sheet_name_fit, index=False)
            print(f"   - Data for {amp:.2f}T written to Excel sheets.")
            
    print(f"\n■ Excel file saved to: {excel_output_path}")

    output_ppt_path = os.path.join(OUTPUT_BASE_PATH, f"Hysteresis_Analysis_{MAT_NAME}.pptx")
    try:
        prs.save(output_ppt_path)
        print(f"\n■ PowerPoint presentation saved to: {output_ppt_path}")
    except Exception as e:
        print(f"\n■ ERROR: Could not save PowerPoint file: {e}")

if __name__ == "__main__":
    # Matplotlibのフォント設定はメイン実行ブロック内が推奨
    plt.rcParams["font.family"] = "Times New Roman"
    plt.rcParams["mathtext.fontset"] = "stix"
    plt.rcParams["font.size"] = 15
    plt.rcParams['figure.dpi'] = 100
    
    main()
    print("\n■ All processes completed.")